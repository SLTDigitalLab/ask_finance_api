from fastapi import APIRouter, Form, HTTPException, Depends
from fastapi.security import HTTPBearer
from pydantic import BaseModel, HttpUrl
from typing import Dict, List, Optional
import logging
import os
import aiohttp
from dotenv import load_dotenv
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
import trafilatura
from readability import Document
from db.psql_connector import DB, default_config
from .vectorstore import *
from langchain_community.tools.tavily_search import TavilySearchResults
from langchain.text_splitter import RecursiveCharacterTextSplitter
from .app_types import AgentState
from .auth import verify_token, token_manager 
import pdfplumber
import docx
import io
import requests
from .context_manager import context_manager
from fastapi import Query
import email
from email import policy
from email.parser import BytesParser
from PIL import Image
import asyncio
import tempfile
import base64
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from google import genai
from google.genai import types




GRAPH_URL = "https://graph.microsoft.com/v1.0"

GEMINI_OCR_URL = "https://api.gemini.ai/v1/vision/ocr"  

logger = logging.getLogger(__name__)

load_dotenv()

TAVILY_API_KEY = os.getenv("TAVILY_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

if TAVILY_API_KEY:
    os.environ["TAVILY_API_KEY"] = TAVILY_API_KEY

if GOOGLE_API_KEY:
    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY




logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

router = APIRouter()

REQUEST_TIMEOUT = 30

security = HTTPBearer()

class ChatRequest(BaseModel):
    query: str
    answer: str = ""
    domain: str = "default"
    chat_mode: str = "short"
    cache_mode: bool = False

class ChatResponse(BaseModel):
    answer: str
    sources: List[str] = []
    pages: List[int] = []
    images: List[str] = []
    chat_id: str
    reasoning_chain: List[str] = []
    domain: str 

class DomainRequest(BaseModel):
    domain: str
    description: Optional[str] = None

class DomainResponse(BaseModel):
    domain: str
    collection_name: str
    status: str
    points_count: Optional[int] = None

class HRKBRequest(BaseModel):
    folder_id: Optional[str] = None
    share_id: Optional[str] = None
    token: str
    domain: str = "hr"
    chunk_size: Optional[int] = 1000
    chunk_overlap: Optional[int] = 200

class LinkRequest(BaseModel):
    urls: List[HttpUrl]
    domain: str 
    chunk_size: Optional[int] = 1000
    chunk_overlap: Optional[int] = 200

class LinkResponse(BaseModel):
    success: bool
    message: str
    document_id: str
    domain: str  
    title: str
    content_length: int
    chunks_created: int
    images: List[str] = []
    internal_links: List[str] = []
    metadata: Dict = {}

class BulkLinkRequest(BaseModel):
    urls: List[HttpUrl]
    domain: str  
    chunk_size: Optional[int] = 1000
    chunk_overlap: Optional[int] = 200
    extract_images: Optional[bool] = False
    extract_links: Optional[bool] = False

class BulkLinkResponse(BaseModel):
    success: bool
    processed: int
    failed: int
    results: List[LinkResponse]
    errors: List[str] = []
    domain: str  

chat_sessions: Dict[str, Dict] = {}
document_collections: Dict[str, Dict] = {}
collection_documents: Dict[str, List[Dict]] = {}

search_tool = TavilySearchResults()

logger = logging.getLogger(__name__)


# Gemini OCR function
os.environ["GENAI_API_KEY"] = os.getenv("GOOGLE_API_KEY")

async def gemini_ocr(file_bytes: bytes, mime_type: str = "image/jpeg") -> str:
    """
    Extract text from an image using Gemini API asynchronously.
    If mime_type is not provided, defaults to 'image/jpeg'.
    """
    try:
        #client = genai.Client()
        client = genai.Client(api_key=GOOGLE_API_KEY)


        # Run synchronous Gemini client in thread executor
        loop = asyncio.get_event_loop()
        response = await loop.run_in_executor(
            None,
            lambda: client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[
                    types.Part.from_bytes(data=file_bytes, mime_type=mime_type),
                    "Extract all text accurately from this flyer/image."
                ]
            )
        )

        return response.text or ""

    except Exception as e:
        print(f"[gemini_ocr] OCR failed: {e}")
        return ""


# EML parser (email)
def build_email_document_structure(msg, body_text: str, attachment_text: str) -> str:
    from_addr = msg.get("From", "")
    to_addr = msg.get("To", "")
    subject = msg.get("Subject", "")
    sent_on = msg.get("Date", "")

    final_text = f"""
From: {from_addr}
Sent on: {sent_on}
To: {to_addr}
Subject: {subject}

Body:
{body_text}

Flyer Text:
{attachment_text}
""".strip()

    return final_text

async def parse_eml(file_bytes: bytes) -> str:
    msg = BytesParser(policy=policy.default).parsebytes(file_bytes)

    body = ""
    attachment_texts = []
    ocr_tasks = []  # list of tuples: (key, task)
    cid_map = {}  # cid -> OCR text

    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            disp = part.get_content_disposition()
            cid = part.get("Content-ID", "").strip("<>")

            # ------------------------
            # BODY TEXT
            # ------------------------
            if ctype == "text/plain" and disp != "attachment":
                body += part.get_content() + "\n"
            elif ctype == "text/html" and disp != "attachment":
                soup = BeautifulSoup(part.get_content(), "html.parser")
                body += soup.get_text() + "\n"

            # ------------------------
            # INLINE IMAGES (cid)
            # ------------------------
            if ctype.startswith("image/") and cid:
                data = part.get_content()
                mime_type = ctype
                task = asyncio.create_task(gemini_ocr(data, mime_type))
                ocr_tasks.append((cid, task))

            # ------------------------
            # ATTACHMENTS
            # ------------------------
            if disp == "attachment":
                filename = (part.get_filename() or "").lower()
                data = part.get_content()

                if filename.endswith(".pdf"):
                    attachment_texts.append(parse_pdf(data))
                elif filename.endswith(".docx"):
                    attachment_texts.append(parse_docx(data))
                elif filename.endswith((".png", ".jpg", ".jpeg")):
                    task = asyncio.create_task(gemini_ocr(data, part.get_content_type()))
                    ocr_tasks.append((filename, task))

    # Run OCR asynchronously
    for key, task in ocr_tasks:
        text = await task
        cid_map[key] = text

    # Replace cids in body with OCR text
    for cid, text in cid_map.items():
        body = body.replace(f"[cid:{cid}]", text)

    # Combine any remaining OCRed attachments that weren’t inline
    attachment_texts.extend([v for k, v in cid_map.items() if k not in body])

    return build_email_document_structure(
        msg,
        body,
        "\n".join(attachment_texts)
    )


# DOCX parser
def parse_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX file with error handling."""
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except Exception as e:
        logger.warning(f"DOCX parsing error: {e}")
        return ""

# PDF parser
def parse_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF file with error handling."""
    text = ""
    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"
    except Exception as e:
        logger.warning(f"PDF parsing error: {e}")
        # Return empty string - caller will skip this file
        return ""
    
    return text.strip()

def parse_text(file_bytes: bytes) -> str:
    """Extract text from text files with error handling."""
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        logger.warning(f"Text parsing error: {e}")
        return ""

# Image parser (OCR)
async def parse_image(file_bytes: bytes) -> str:
    try:
        Image.open(io.BytesIO(file_bytes))  # Validate image
        ocr_text = await gemini_ocr(file_bytes)
        return f"[IMAGE OCR] {ocr_text.strip()}" if ocr_text else ""
    except Exception as e:
        print(f"Image parsing failed: {e}")
        return ""

# PPTX parser with OCR for images
async def parse_pptx_smart(file_bytes: bytes) -> str:
    """
    Fully enhanced PPTX extractor:
    - Slide text
    - SmartArt text
    - Grouped shape text
    - Table text
    - Speaker notes
    - Images (ppt/media, embeddings, charts) → OCR via Gemini
    """

    prs = Presentation(io.BytesIO(file_bytes))
    extracted_text = []

    # ------------------------------------------
    # Helper: recursively extract text from shapes
    # ------------------------------------------
    def extract_shape_text(shape):
        try:
            # Basic text
            if hasattr(shape, "text") and shape.text.strip():
                extracted_text.append(shape.text.strip())

            # Text frame
            if hasattr(shape, "text_frame") and shape.text_frame:
                for p in shape.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        extracted_text.append(t)

            # Table text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            extracted_text.append(t)

            # Grouped shapes / SmartArt
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in shape.shapes:
                    extract_shape_text(sub)

        except Exception:
            pass

    # ------------------------------------------
    # Extract text from slides + notes
    # ------------------------------------------
    for slide in prs.slides:
        for shape in slide.shapes:
            extract_shape_text(shape)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                extracted_text.append(f"[SPEAKER NOTES] {notes}")

    # ------------------------------------------
    # Extract images from PPTX ZIP
    # ------------------------------------------
    pptx_zip = zipfile.ZipFile(io.BytesIO(file_bytes))
    image_extensions = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff")

    image_files = [
        fn for fn in pptx_zip.namelist()
        if fn.lower().endswith(image_extensions)
    ]

    # Also check ppt/media, ppt/embeddings
    for fn in pptx_zip.namelist():
        lower_fn = fn.lower()
        if (
            lower_fn.startswith("ppt/media/")
            or lower_fn.startswith("ppt/embeddings/")
        ):
            if lower_fn.endswith(image_extensions):
                if fn not in image_files:
                    image_files.append(fn)

    # ------------------------------------------
    # OCR each image
    # ------------------------------------------
    for img_name in image_files:
        try:
            img_bytes = pptx_zip.read(img_name)

            # Validate correct image type
            try:
                Image.open(io.BytesIO(img_bytes))
            except:
                continue  # skip broken or non-image files

            # OCR using Gemini
            ocr_text = await gemini_ocr(img_bytes)

            if ocr_text and ocr_text.strip():
                extracted_text.append(f"[IMAGE OCR] {ocr_text.strip()}")

        except Exception as e:
            print(f"Failed OCR for image {img_name}: {e}")

    # ------------------------------------------
    # Clean + dedupe
    # ------------------------------------------
    cleaned = []
    seen = set()

    for line in extracted_text:
        line = line.strip()
        if line and line not in seen:
            cleaned.append(line)
            seen.add(line)

    return "\n".join(cleaned)


def get_chat_history(chat_id: str, domain: Optional[str] = None):
    """Get chat history from database, optionally filtered by domain."""
    db = None
    try:
        db = DB(default_config())
        cursor = db.conn.cursor()
        
        if domain:
            query = """
            SELECT role, message FROM ask_hr_history
            WHERE chat_id = %s AND domain = %s
            ORDER BY timestamp ASC
            """
            cursor.execute(query, (chat_id, domain))
        else:
            query = """
            SELECT role, message FROM ask_hr_history
            WHERE chat_id = %s
            ORDER BY timestamp ASC
            """
            cursor.execute(query, (chat_id,))
        
        result = cursor.fetchall()
        
        history = []
        for row in result:
            history.append({
                "role": row[0],
                "content": row[1]
            })
        
        return history
        
    except Exception as e:
        logger.error(f"Error loading chat history: {e}")
        return []
    finally:
        if db:
            try:
                db.close()
            except:
                pass

async def fetch_page_content(session: aiohttp.ClientSession, url: str) -> tuple[str, Dict]:
    """Fetch and extract content from a web page."""
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }
    
    try:
        async with session.get(str(url), headers=headers, timeout=aiohttp.ClientTimeout(total=30)) as response:
            if response.status != 200:
                raise HTTPException(status_code=400, detail=f"Failed to fetch URL: HTTP {response.status}")
            
            html_content = await response.text()

            content = ""
            title = ""
            metadata = {}
            
            try:
                extracted = trafilatura.extract(html_content, include_comments=False, include_tables=True)
                if extracted:
                    content = extracted

                    metadata_extracted = trafilatura.extract_metadata(html_content)
                    if metadata_extracted:
                        title = metadata_extracted.title or ""
                        metadata.update({
                            'author': metadata_extracted.author,
                            'date': str(metadata_extracted.date) if metadata_extracted.date else None,
                            'description': metadata_extracted.description,
                            'categories': metadata_extracted.categories,
                            'tags': metadata_extracted.tags
                        })
            except Exception as e:
                logger.warning(f"Trafilatura extraction failed: {e}")
            
            if not content:
                try:
                    doc = Document(html_content)
                    content = doc.summary()
                    title = doc.title()
                except Exception as e:
                    logger.warning(f"Readability extraction failed: {e}")
            
            if not content:
                try:
                    soup = BeautifulSoup(html_content, 'html.parser')
                    
                    for script in soup(["script", "style"]):
                        script.decompose()
                    
                    if not title:
                        title_tag = soup.find('title')
                        title = title_tag.get_text().strip() if title_tag else ""
                    
                    content = soup.get_text()
                    
                    lines = (line.strip() for line in content.splitlines())
                    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                    content = ' '.join(chunk for chunk in chunks if chunk)
                    
                except Exception as e:
                    logger.error(f"BeautifulSoup extraction failed: {e}")
                    content = html_content
            
            if not content:
                raise HTTPException(status_code=400, detail="Failed to extract content from the webpage")
            
            metadata.update({
                'url': str(url),
                'title': title,
                'content_length': len(content),
                'extraction_method': 'trafilatura' if 'trafilatura' in str(type(content)) else 'readability' if 'readability' in str(type(content)) else 'beautifulsoup'
            })
            
            return content, metadata
            
    except aiohttp.ClientError as e:
        logger.error(f"Network error fetching {url}: {e}")
        raise HTTPException(status_code=400, detail=f"Network error: {str(e)}")
    except Exception as e:
        logger.error(f"Error processing {url}: {e}")
        raise HTTPException(status_code=500, detail=f"Processing error: {str(e)}")

async def fetch_onedrive_folder_docs(folder_id: str = None, token: str = None, share_id: str = None):
    """Fetch all files in a OneDrive folder and extract their text."""
    
    if share_id:
        url = f"{GRAPH_URL}/shares/{share_id}/driveItem/children"
    elif folder_id:
        url = f"{GRAPH_URL}/me/drive/items/{folder_id}/children"
    else:
        raise ValueError("Either folder_id or share_id must be provided")
    
    headers = {"Authorization": f"Bearer {token}"}

    response = requests.get(url, headers=headers, timeout=60)
    response.raise_for_status()
    items = response.json().get("value", [])
    
    logger.info(f"Found {len(items)} items to process")
    
    # Log all files that will be processed
    if items:
        logger.info("Files to process:")
        for i, item in enumerate(items):
            logger.info(f"  {i+1}: {item.get('name', 'Unknown')}")
    
    docs = []
    processed_count = 0
    error_count = 0
    
    for item in items:
        try:
            item_name = item.get("name", "Unknown")
            logger.info(f"Processing: {item_name}")
            
            if "@microsoft.graph.downloadUrl" not in item:
                logger.warning(f"Skipping {item_name} - no download URL")
                error_count += 1
                continue
            
            download_url = item["@microsoft.graph.downloadUrl"]
            logger.info(f"  Downloading from: {download_url[:50]}...")
            
            file_resp = requests.get(download_url, timeout=30)
            file_resp.raise_for_status()
            file_bytes = file_resp.content
            name = item["name"].lower()

            if name.endswith(".docx"):
                file_text = parse_docx(file_bytes)
            elif name.endswith(".pdf"):
                file_text = parse_pdf(file_bytes)
            elif name.endswith(".eml"):
                file_text = await parse_eml(file_bytes)  # New support for emails
            elif name.endswith(".pptx"):
                file_text = await parse_pptx_smart(file_bytes)
            elif name.endswith((".jpg", ".jpeg", ".png")):
                file_text = await gemini_ocr(file_bytes)
            elif name.endswith((".txt", ".md", ".json")):
                file_text = parse_text(file_bytes)
            else:
                logger.warning(f"Skipping unsupported file type: {item['name']}")
                continue

            if not file_text or not file_text.strip():
                logger.warning(f"Skipping {item_name} - empty content")
                error_count += 1
                continue
            
            logger.info(f"  Extracted {len(file_text)} characters")
            
            # Log a preview of the content
            if file_text:
                preview = file_text[:100].replace('\n', ' ')
                logger.info(f"  Content preview: {preview}...")

            docs.append({
                "id": item["id"],
                "name": item_name,
                "content": file_text,
                "web_url": item.get("webUrl"),
                "share_id": item.get("shareId"),
            })
            processed_count += 1
            logger.info(f"✅ Successfully processed: {item_name}")
            
        except Exception as e:
            logger.error(f"❌ Error processing {item.get('name', 'Unknown')}: {e}")
            error_count += 1
            continue

    logger.info(f"Processing complete: {processed_count} successful, {error_count} failed")
    
    if processed_count == 0 and error_count > 0:
        logger.warning("No files were successfully processed!")
    
    return docs

def chunk_content(content: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
    """Split content into chunks for processing."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )
    
    chunks = text_splitter.split_text(content)
    return chunks

@router.post("/domains/create", tags=["Domains"], response_model=DomainResponse)
async def create_domain(
    request: DomainRequest,
    token: str = Depends(token_manager.verify_admin_token)
):
    """Create a new domain collection."""
    try:
        status = create_collection(domain=request.domain)
        stats = get_domain_stats(request.domain)
        
        return DomainResponse(
            domain=request.domain,
            collection_name=stats.get("collection_name", ""),
            status=status,
            points_count=stats.get("points_count", 0)
        )
    except Exception as e:
        logger.error(f"Error creating domain: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/domains", tags=["Domains"])
async def list_domains(
    token: str = Depends(verify_token)
):
    """List all available domains."""
    try:
        collections = get_all_collections()
        
        domain_info = []
        for col in collections:
            stats = get_domain_stats(col["domain"])
            domain_info.append({
                "domain": col["domain"],
                "collection_name": col["collection_name"],
                "points_count": stats.get("points_count", 0),
                "status": stats.get("status", "unknown")
            })
        
        return {
            "domains": domain_info,
            "total_count": len(domain_info)
        }
    except Exception as e:
        logger.error(f"Error listing domains: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/domains/{domain}/stats", tags=["Domains"])
async def get_domain_statistics(
    domain: str,
    token: str = Depends(verify_token)
):
    """Get statistics for a specific domain."""
    try:
        stats = get_domain_stats(domain)
        if not stats.get("exists"):
            raise HTTPException(status_code=404, detail=f"Domain '{domain}' not found")
        return stats
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error getting domain stats: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/domains/{domain}", tags=["Domains"])
async def delete_domain(
    domain: str,
    token: str = Depends(token_manager.verify_admin_token)
):
    """Delete a domain and its collection."""
    try:
        delete_collection(domain=domain)
        return {
            "status": "success",
            "message": f"Domain '{domain}' deleted successfully"
        }
    except Exception as e:
        logger.error(f"Error deleting domain: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/add_data", tags=["Vectorstore"])
async def add_data_to_collection(
    request: LinkRequest,
):
    """Add web content to a specific domain."""
    result = {"status": "", "message": "", "domain": request.domain}

    try:
        async with aiohttp.ClientSession() as session:
            contents = ""
            for link in request.urls:
                content, metadata = await fetch_page_content(session, str(link))
                contents += content

            chunks = chunk_content(contents, request.chunk_size, request.chunk_overlap)

            logger.info(f"Successfully processed {len(request.urls)} links into {len(chunks)} chunks for domain '{request.domain}'")

            upsert_result = add_texts(chunks, domain=request.domain)

            result["status"] = "success"
            result["message"] = f"Added {len(chunks)} chunks to domain '{request.domain}'"

            return result

    except Exception as e:
        logger.error(f"Error processing links for domain '{request.domain}': {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))



@router.post("/add_hr_kb", tags=["Vectorstore"])
async def add_hr_kb_to_collection(
    request: HRKBRequest,
):
    """Add OneDrive documents to a specific domain with OneDrive IDs."""
    result = {"status": "", "message": "", "domain": request.domain}

    try:
        docs = await fetch_onedrive_folder_docs(request.folder_id, request.token, share_id=request.share_id)
        
        if not docs:
            result["status"] = "failed"
            result["message"] = "No documents found in OneDrive folder"
            return result

        all_chunks = []
        document_names = []
        onedrive_ids = []
        onedrive_urls = []
        
        for doc in docs:
            chunks = chunk_content(doc["content"], request.chunk_size, request.chunk_overlap)
            all_chunks.extend(chunks)
            
            # Add document metadata for each chunk
            document_names.extend([doc["name"]] * len(chunks))
            onedrive_ids.extend([doc["id"]] * len(chunks))
            
            # Try to get the web URL, fallback to constructing one
            if doc.get("web_url"):
                onedrive_urls.extend([doc["web_url"]] * len(chunks))
            else:
                # Construct a web URL from the ID
                web_url = f"https://onedrive.live.com/?id={doc['id']}&cid=7E1063C6DE897DDC"
                onedrive_urls.extend([web_url] * len(chunks))
            
            logger.info(f"Processed OneDrive file {doc['name']} (ID: {doc['id']}) into {len(chunks)} chunks")

        # Pass OneDrive IDs and URLs to add_texts
        upsert_result = add_texts(
            all_chunks, 
            domain=request.domain,
            document_names=document_names,
            onedrive_ids=onedrive_ids,
            onedrive_urls=onedrive_urls
        )

        result["status"] = "success"
        result["message"] = f"Inserted {len(all_chunks)} chunks from {len(docs)} OneDrive docs to domain '{request.domain}'"
        return result
    
    
    except Exception as e:
        logger.error(f"Error processing OneDrive folder for domain '{request.domain}': {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))



@router.get("/chunks/{domain}", tags=["Vectorstore"])
async def list_domain_chunks(
    domain: str,
    token: str = Depends(verify_token)
):
    """Get all chunks in a specific domain."""
    try:
        raw_domain = domain
        domain = (domain or "").strip().lower()

        # Keep consistent with base.py
        DOMAIN_ALIASES = {
            "ask_fiannce": "ask_finance",
            "ask_fianance": "ask_finance",
        }
        domain = DOMAIN_ALIASES.get(domain, domain)

        if raw_domain != domain:
            logger.warning(f"[CHUNKS] Domain normalized: '{raw_domain}' -> '{domain}'")

        chunks = get_all_points(domain=domain)
        return {
            "domain": domain,
            "chunks": chunks,
            "total_count": len(chunks)
        }
    except Exception as e:
        logger.error(f"Error listing chunks for domain '{domain}': {e}")
        raise HTTPException(status_code=500, detail=str(e))



def document_search_agent(state: Dict[str, Any]) -> Dict[str, Any]:
    """
    Document search with intelligent caching for ANY follow-up question and document name extraction for reference preview.
    
    Caching Strategy:
    - Level 1: Cache retrieved document content (for questions about "that document")
    - Level 2: Cache last assistant answer (for questions about "explain above")
    """
    query = state["query"]
    original_query = state.get("original_query", query)
    chat_id = state.get("chat_id")
    domain = state.get("collection_id")
    
    logger.info(f"[DOC_AGENT] Query: '{original_query}'")
    logger.info(f"[DOC_AGENT] Domain: '{domain}'")

    # Check if this is a simple greeting that shouldn't have references
    is_simple_greeting = any(
        greeting in original_query.lower() 
        for greeting in ["hello", "hi", "hey", "good morning", "good afternoon", "good evening", "how are you"]
    )

    # Check for personal info queries
    personal_info_keywords = [
        "my name", "your name", "who am i", "who are you", "i am", "my age",
        "my birthday", "my details", "personal information", "about me"
    ]
    
    is_personal_info_query = any(
        keyword in original_query.lower() 
        for keyword in personal_info_keywords
    )
    
    # Clear any previous references for simple greetings or personal info
    if is_simple_greeting or is_personal_info_query:
        logger.info(f"[DOC_AGENT] {'Simple greeting' if is_simple_greeting else 'Personal info query'} detected - clearing all document references")
        state["document_context"] = ""
        state["document_names"] = []
        state["onedrive_urls"] = []
        state["has_reference"] = False
        state["document_found"] = False
        state["reasoning_chain"].append(f"Document Search: {'Simple greeting' if is_simple_greeting else 'Personal info query'} - no document search needed")
        return state

    is_vague = context_manager.is_followup_question(original_query or query)
    
    # Get cached data
    last_doc_content = context_manager.get_context("document_agent", "last_document_content", chat_id)
    last_answer = context_manager.get_last_answer_snippet(chat_id, max_length=1000)
    
    logger.info(f"[DOC_AGENT] Is vague: {is_vague}")
    logger.info(f"[DOC_AGENT] Has cached doc: {last_doc_content is not None}")
    
    followup_type = None
    
    if is_vague:
        query_lower = original_query.lower()
        
        document_refs = [
            "that document", "this document", "the document",
            "what does it say", "what does it contain", "what's in it",
            "document content", "document policy"
        ]
        
        if any(ref in query_lower for ref in document_refs):
            followup_type = "DOCUMENT_SPECIFIC"
            logger.info(f"[DOC_AGENT] Follow-up type: DOCUMENT_SPECIFIC")
        
        elif any(word in query_lower for word in ['explain', 'elaborate', 'detail', 'more about']):
            answer_refs = [
                'explain all', 'explain those', 'explain these', 'explain them',
                'explain above', 'above thing', 'above mentioned',
                'tell me more', 'more details', 'elaborate',
                'break down', 'go deeper'
            ]
            
            if any(ref in query_lower for ref in answer_refs):
                followup_type = "ANSWER_ELABORATION"
                logger.info(f"[DOC_AGENT] Follow-up type: ANSWER_ELABORATION")
        
        elif len(original_query.split()) <= 5:
            if any(word in query_lower for word in ['it', 'them', 'those', 'these']):
                followup_type = "GENERIC_VAGUE"
                logger.info(f"[DOC_AGENT] Follow-up type: GENERIC_VAGUE")
    
    if followup_type == "DOCUMENT_SPECIFIC" and last_doc_content:
        logger.info(f"[DOC_AGENT] ✓ Strategy 1: Reusing cached document content")

        # Get cached document names and URLs
        cached_names = context_manager.get_context("document_agent", "last_document_names", chat_id) or []
        cached_urls = context_manager.get_context("document_agent", "last_onedrive_urls", chat_id) or []

        state["document_context"] = last_doc_content
        state["document_found"] = True
        state["sources"] = context_manager.get_context("document_agent", "last_document_ids", chat_id) or []
        state["document_names"] = cached_names
        state["onedrive_urls"] = cached_urls
        state["has_reference"] = len(cached_names) > 0
        state["reasoning_chain"].append("Document Search: Reused cached document (document-specific follow-up)")
        
        context_manager.update_entity_memory(chat_id, query, role="user")
        return state
    
    if followup_type == "ANSWER_ELABORATION" and last_doc_content:
        logger.info(f"[DOC_AGENT] ✓ Strategy 2: Reusing document for answer elaboration")

        # Get cached document names and URLs
        cached_names = context_manager.get_context("document_agent", "last_document_names", chat_id) or []
        cached_urls = context_manager.get_context("document_agent", "last_onedrive_urls", chat_id) or []
        
        # Add the previous answer as additional context
        combined_context = last_doc_content
        
        if last_answer:
            combined_context = f"Previous Answer:\n{last_answer}\n\n---\n\nOriginal Document Content:\n{last_doc_content}"
        
        state["document_context"] = combined_context
        state["document_found"] = True
        state["sources"] = context_manager.get_context("document_agent", "last_document_ids", chat_id) or []
        state["document_names"] = cached_names
        state["onedrive_urls"] = cached_urls
        state["has_reference"] = len(cached_names) > 0
        state["reasoning_chain"].append("Document Search: Reused document + previous answer (elaboration request)")
        
        context_manager.update_entity_memory(chat_id, query, role="user")
        return state
    
    if followup_type == "GENERIC_VAGUE" and last_doc_content:
        logger.info(f"[DOC_AGENT] ✓ Strategy 3: Attempting to reuse for generic vague query")
        
        # Get cached document names and URLs
        cached_names = context_manager.get_context("document_agent", "last_document_names", chat_id) or []
        cached_urls = context_manager.get_context("document_agent", "last_onedrive_urls", chat_id) or []
        
        state["document_context"] = last_doc_content
        state["document_found"] = True
        state["sources"] = context_manager.get_context("document_agent", "last_document_ids", chat_id) or []
        state["document_names"] = cached_names
        state["onedrive_urls"] = cached_urls
        state["has_reference"] = len(cached_names) > 0
        state["reasoning_chain"].append("Document Search: Reused cached document (generic follow-up)")
        
        context_manager.update_entity_memory(chat_id, query, role="user")
        return state
    
    logger.info(f"[DOC_AGENT] Strategy 4: Performing new document search")
    
    # Extract entities
    entities = context_manager.extract_entities(query)
    doc_codes = entities.get("CUSTOM", [])
    
    if not doc_codes and is_vague:
        recent_entities = context_manager.get_recent_entities(chat_id, limit=5)
        doc_codes = [e for e in recent_entities if any(char.isdigit() for char in e)]
        logger.info(f"[DOC_AGENT] Document codes from context: {doc_codes}")
    
    # Multi-strategy search
    all_docs = []
    
    if doc_codes:
        for doc_code in doc_codes[:2]:
            docs = search_similar(doc_code, domain=domain, limit=5)
            if docs:
                all_docs.extend(docs)
    
    if len(all_docs) == 0:
        docs = search_similar(original_query, domain=domain, limit=5)
        if docs:
            all_docs.extend(docs)
    
    if len(all_docs) == 0:
        # Try case-insensitive search for document names
        # Extract potential document names from query
        query_words = original_query.lower().split()
        potential_doc_names = []
        
        # Look for phrases that might be document titles
        for i in range(len(query_words)):
            for j in range(i+1, min(i+5, len(query_words))):  # Look for 2-5 word phrases
                phrase = " ".join(query_words[i:j+1])
                # Check if this phrase might be a document title
                if len(phrase) > 3 and not phrase in ["the", "and", "for", "with", "about"]:
                    potential_doc_names.append(phrase)
        
        # Search for each potential document name
        for doc_name in potential_doc_names[:3]:  # Limit to top 3
            docs = search_similar(doc_name, domain=domain, limit=3)
            if docs:
                all_docs.extend(docs)
                logger.info(f"[DOC_AGENT] Found documents using phrase: '{doc_name}'")
    
    logger.info(f"[DOC_AGENT] Search returned {len(all_docs)} documents")
    
    if not all_docs:
        logger.error("[DOC_AGENT] ERROR: No documents found in search!")
        state["document_found"] = False
        state["has_reference"] = False
        state["document_names"] = []
        state["onedrive_urls"] = []
        state["reasoning_chain"].append("Document Search: No documents found")
        return state
    
    # Deduplicate and sort
    seen = set()
    unique_docs = []
    for doc in all_docs:
        doc_id = doc.get('id')
        if doc_id not in seen:
            seen.add(doc_id)
            unique_docs.append(doc)
    
    unique_docs.sort(key=lambda x: x.get('score', 0), reverse=True)
    
    # Log detailed info about top docs
    logger.info(f"[DOC_AGENT] Top {len(unique_docs[:5])} unique documents:")
    for i, doc in enumerate(unique_docs[:5], 1):
        score = doc.get('score', 0)
        content = doc.get('payload', {}).get('page_content', '')
        doc_name = doc.get('payload', {}).get('document_name', 'Unknown')
        onedrive_url = doc.get('payload', {}).get('onedrive_url', 'No URL')
        logger.info(f"  [{i}] Score: {score:.3f}")
        logger.info(f"      Doc: {doc_name}")
        logger.info(f"      URL: {onedrive_url}")
        logger.info(f"      Content preview: {content[:100]}...")
        logger.info(f"      Content length: {len(content)} chars")
    
    # Build context
    top_docs = unique_docs[:5]
    raw_contexts = []
    document_names = []
    onedrive_urls = []
    doc_ids = []
    
    logger.info(f"[DOC_AGENT] Processing {len(top_docs)} top documents...")
    
    for i, doc in enumerate(top_docs):
        content = doc.get('payload', {}).get('page_content', '')
        doc_name = doc.get('payload', {}).get('document_name', '')
        onedrive_url = doc.get('payload', {}).get('onedrive_url', '')
        
        logger.info(f"[DOC_AGENT] Doc {i+1}: name='{doc_name}', content_type={type(content)}, content_len={len(content) if content else 0}")
        
        # Simple check - just make sure content exists
        if content and isinstance(content, str) and content.strip():
            raw_contexts.append(content.strip())
            logger.info(f"[DOC_AGENT]   Added content: {len(content.strip())} chars")
            
            if doc_name and isinstance(doc_name, str) and doc_name not in document_names:
                document_names.append(doc_name)
                logger.info(f"[DOC_AGENT]   Added document name: {doc_name}")
            
            if onedrive_url and isinstance(onedrive_url, str) and onedrive_url not in onedrive_urls:
                onedrive_urls.append(onedrive_url)
                logger.info(f"[DOC_AGENT]   Added OneDrive URL: {onedrive_url}")
            
            doc_id = doc.get('id')
            if doc_id and doc_id not in doc_ids:
                doc_ids.append(doc_id)
        else:
            logger.warning(f"[DOC_AGENT]   Skipping invalid content: type={type(content)}, length={len(content) if content else 0}")

    logger.info(f"[DOC_AGENT] Extracted {len(raw_contexts)} raw contexts")
    logger.info(f"[DOC_AGENT] Document names found: {document_names}")
    logger.info(f"[DOC_AGENT] OneDrive URLs found: {onedrive_urls}")

    if not raw_contexts:
        logger.error("[DOC_AGENT] ERROR: No valid content extracted!")
        state["document_found"] = False
        state["has_reference"] = False
        state["document_names"] = []
        state["onedrive_urls"] = []
        state["answer"] = "I found documents but couldn't extract readable content."
        return state
    
    # Check if documents are actually relevant to the query
    relevant_contexts = []
    relevant_doc_names = []
    relevant_urls = []
    relevant_doc_ids = []
    
    for i, content in enumerate(raw_contexts):
        # Try to get the corresponding document name and URL
        doc_name = ""
        if i < len(document_names):
            doc_name = document_names[i]
        
        onedrive_url = ""
        if i < len(onedrive_urls):
            onedrive_url = onedrive_urls[i]
        
        doc_id = ""
        if i < len(doc_ids):
            doc_id = doc_ids[i]
        
        # Check if this document is relevant to the query
        if is_document_relevant_to_query(content, original_query):
            relevant_contexts.append(content.strip())
            if doc_name and doc_name not in relevant_doc_names:
                relevant_doc_names.append(doc_name)
            if onedrive_url and onedrive_url not in relevant_urls:
                relevant_urls.append(onedrive_url)
            if doc_id and doc_id not in relevant_doc_ids:
                relevant_doc_ids.append(doc_id)
            logger.info(f"[DOC_AGENT] Document '{doc_name}' deemed RELEVANT to query")
        else:
            logger.info(f"[DOC_AGENT] Document '{doc_name}' deemed NOT RELEVANT to query")
    
    if not relevant_contexts:
        logger.info("[DOC_AGENT] No documents found relevant to the query")
        state["document_found"] = False
        state["has_reference"] = False
        state["document_names"] = []
        state["onedrive_urls"] = []
        state["reasoning_chain"].append("Document Search: No relevant documents found")
        return state
    
    # Use only relevant documents
    context = "\n\n".join(relevant_contexts)
    document_names = relevant_doc_names
    onedrive_urls = relevant_urls
    doc_ids = relevant_doc_ids
    
    logger.info(f"[DOC_AGENT] After relevance filtering: {len(relevant_contexts)} relevant contexts, {len(document_names)} document names")

    context = "\n\n".join(raw_contexts)
    logger.info(f"[DOC_AGENT] Built context (length: {len(context)})")
    logger.info(f"[DOC_AGENT] Context first 500 chars: {context[:500]}...")
    
    # Cache for future follow-ups
    context_manager.update_context("document_agent", "last_document_ids", doc_ids, chat_id)
    context_manager.update_context("document_agent", "last_document_names", document_names, chat_id)
    context_manager.update_context("document_agent", "last_onedrive_urls", onedrive_urls, chat_id)
    context_manager.update_context("document_agent", "last_document_content", context, chat_id)
    context_manager.update_context("document_agent", "last_query", original_query, chat_id)
    
    logger.info(f"[DOC_AGENT] ✓ Cached {len(top_docs)} documents with names: {document_names}")
    
    # Update state
    state["document_context"] = context.strip()
    state["document_found"] = True
    state["sources"] = doc_ids
    state["document_names"] = document_names
    state["onedrive_urls"] = onedrive_urls
    state["has_reference"] = len(document_names) > 0
    
    # Log what we're setting in state
    logger.info(f"[DOC_AGENT] Setting in state:")
    logger.info(f"  - document_context: {len(state['document_context'])} chars")
    logger.info(f"  - document_names: {state['document_names']}")
    logger.info(f"  - onedrive_urls: {state['onedrive_urls']}")
    logger.info(f"  - has_reference: {state['has_reference']}")
    
    state["reasoning_chain"].append(f"Document Search: Retrieved {len(top_docs)} new documents")
    
    # Update entity memory
    context_manager.update_entity_memory(chat_id, query, role="user")
    context_manager.update_entity_memory(chat_id, context[:500], role="assistant")
    
    if domain:
        context_manager.update_context("document_agent", "last_domain", domain, chat_id)
    
    return state


def is_document_relevant_to_query(document_content: str, query: str) -> bool:
    """
    Check if document content is actually relevant to the query.
    Returns True if the document appears to contain information related to the query.
    """
    query_lower = query.lower()
    content_lower = document_content.lower()
    
    # Extract key terms from query
    query_terms = set(query_lower.split())
    
    # Remove common stop words
    stop_words = {"the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for", 
                  "of", "with", "by", "is", "are", "was", "were", "be", "been", "being",
                  "have", "has", "had", "do", "does", "did", "will", "would", "shall",
                  "should", "may", "might", "must", "can", "could", "what", "where",
                  "when", "why", "how", "who", "which", "this", "that", "these", "those"}
    
    query_terms = query_terms - stop_words
    
    # Check if any query terms appear in the document
    relevant_terms_found = 0
    for term in query_terms:
        if len(term) > 2 and term in content_lower:
            relevant_terms_found += 1
    
    # If at least 50% of non-stopword terms are found, consider it relevant
    if len(query_terms) > 0:
        relevance_ratio = relevant_terms_found / len(query_terms)
        logger.info(f"[DOC_RELEVANCE] Query terms: {query_terms}")
        logger.info(f"[DOC_RELEVANCE] Found {relevant_terms_found}/{len(query_terms)} terms, ratio: {relevance_ratio:.2f}")
        return relevance_ratio >= 0.3
    
    return False